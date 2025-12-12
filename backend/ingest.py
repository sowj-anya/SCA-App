import argparse
import os
from pathlib import Path
from typing import Iterable, List

from langchain_text_splitters import RecursiveCharacterTextSplitter
from pypdf import PdfReader
import docx2txt
from pptx import Presentation

from backend.config import settings
from backend.rag import build_index


def load_text_from_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ""
        pages.append(text)
    return "\n".join(pages)


def load_text_from_docx(path: Path) -> str:
    """Extract text from .docx file"""
    try:
        return docx2txt.process(str(path))
    except Exception as e:
        raise ValueError(f"Error reading DOCX file {path}: {str(e)}")


def load_text_from_pptx(path: Path) -> str:
    """Extract text from .pptx file"""
    try:
        prs = Presentation(str(path))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        raise ValueError(f"Error reading PPTX file {path}: {str(e)}")


def iter_documents(data_dir: Path) -> Iterable[tuple[str, str]]:
    for file in data_dir.glob("**/*"):
        if file.is_dir():
            continue
        suffix = file.suffix.lower()
        if suffix in {".txt", ".md"}:
            yield file.name, file.read_text(encoding="utf-8", errors="ignore")
        elif suffix == ".pdf":
            yield file.name, load_text_from_pdf(file)
        elif suffix in {".docx", ".doc"}:
            yield file.name, load_text_from_docx(file)
        elif suffix in {".pptx", ".ppt"}:
            yield file.name, load_text_from_pptx(file)


def chunk_documents(docs: Iterable[tuple[str, str]]) -> List[dict]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        separators=["\n\n", "\n", ". ", " "],
    )
    chunks: List[dict] = []
    for source, text in docs:
        if not text.strip():
            continue
        for i, chunk in enumerate(splitter.split_text(text)):
            chunks.append(
                {
                    "id": f"{source}_{i}",
                    "text": chunk.strip(),
                    "source": source,
                }
            )
    return chunks


def ingest(data_dir: str | Path = settings.data_dir) -> None:
    data_dir = Path(data_dir)
    if not data_dir.exists():
        raise FileNotFoundError(f"Data directory not found: {data_dir}")
    docs = list(iter_documents(data_dir))
    if not docs:
        raise RuntimeError("No documents found to ingest.")
    chunks = chunk_documents(docs)
    build_index(chunks)
    print(f"Ingestion complete. Chunks: {len(chunks)}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Ingest documents into the vector index.")
    parser.add_argument("--data-dir", default=settings.data_dir, help="Directory containing source documents.")
    args = parser.parse_args()
    ingest(args.data_dir)


if __name__ == "__main__":
    main()


